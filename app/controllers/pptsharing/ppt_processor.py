from typing import cast, Dict, List, Optional, Tuple
import os
import io
import shutil
import tempfile
import subprocess
import threading

from pptx import Presentation as open_presentation
from pptx.presentation import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.util import Pt
from pptx.text.text import TextFrame
from pptx.dml.color import RGBColor
from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes
from pptx.shapes.autoshape import Shape
from PIL import Image
from PIL import ImageDraw, ImageFont

from app.utils.convert_ppt_to_pptx import util_convert_ppt_to_pptx

try:
    import subprocess
    _HAS_LIBREOFFICE = (
        subprocess.run(["which", "libreoffice"], capture_output=True)
    ).returncode == 0
except Exception:
    _HAS_LIBREOFFICE = False

RENDER_QUALITY: int = 85
DEFAULT_WIDTH: int = 1280
DEFAULT_HEIGHT: int = 720

class PPTProcessor:
    def __init__(self) -> None:
        self.loaded_file: Optional[str] = None
        self.total_slides: int = 0
        self.current_slide: int = 0
        self.slide_cache: Optional[List[bytes]] = []

        self.pdf_path:     Optional[str] = None   # path to converted PDF
        self.pdf_tmpdir:   Optional[str] = None   # temp dir holding PDF file converted
        

    def load_file(self, filepath: str) -> bool:
        try:
            if not os.path.isfile(filepath):
                return False

            actual_path = filepath
            if self.pdf_tmpdir:
                shutil.rmtree(self.pdf_tmpdir, ignore_errors=True)
                self.pdf_tmpdir = None
                self.pdf_path   = None
 
            actual_path = filepath
            if filepath.lower().endswith(".ppt"):
                converted = util_convert_ppt_to_pptx(filepath)
                if converted is None:
                    return False
                actual_path = converted


            prs: Presentation = open_presentation(actual_path)
            self.total_slides = len(prs.slides)
            self.current_slide = 0
            self.loaded_file = filepath
            self.slide_cache = [b""] * self.total_slides  # pre-allocate slots

            # Render slides in background so load_file returns immediately
            def render_all() -> None:
                # convert to PDF once (LibreOffice)
                if _HAS_LIBREOFFICE:
                    self.pdf_path, self.pdf_tmpdir = self._convert_to_pdf(actual_path)
 
                # render each slide
                for idx in range(self.total_slides):
                    self.slide_cache[idx] = self._render_slide(prs, idx)
                    print(f"[PPTProcessor] Rendered slide {idx + 1}/{self.total_slides}")


            thread = threading.Thread(target=render_all, daemon=True)
            thread.start()

            print(f"[PPTProcessor] Loaded '{filepath}' ({self.total_slides} slides), rendering in background")
            return True

        except Exception as exc:
            print(f"[PPTProcessor] load_file error: {exc}")
            self.loaded_file = None
            self.total_slides = 0
            self.slide_cache = []
            return False 

    # returns the bytes of the images of a specific slide num
    def extract_slide(self, slide_num: int) -> bytes:
        if not self.slide_cache:
            return b""

        if slide_num < 0 or slide_num >= len(self.slide_cache):
            return b""

        return self.slide_cache[slide_num]

    # getter for the current file total slides count
    def get_total_slides(self) -> int:
        return self.total_slides

    # getter for the current position of the slide
    def get_current_slideno(self) -> int:
        return self.current_slide

    # utility to switch slide to the next one, fails if the slide no already in the end
    def next(self) -> bool:
        if self.total_slides == 0:
            return False

        if self.current_slide >= self.total_slides - 1:
            return False

        self.current_slide += 1
        return True
    
    # the same as the next(), this is the reverse
    def prev(self) -> bool:
        if self.total_slides == 0:
            return False

        if self.current_slide <= 0:
            return False

        self.current_slide -= 1
        return True

    # to jump from a slide to the requested slide no
    def goto_slide(self, slideno: int) -> bool:
        if slideno < 0 or slideno >= self.total_slides:
            return False

        self.current_slide = slideno
        return True

    # internal function to render each of the slide image,
    # if libreoffice is available, will use libreoffice
    # otherwise fallbacks to python-pptx (minimal loader)
    def _render_slide(self, prs: Presentation, slide_idx: int) -> bytes:
        if _HAS_LIBREOFFICE and self.loaded_file:
            result = self._render_via_pdf(slide_idx)
            if result:
                return result

        return self._render_python_fallback(prs=prs, slide_idx=slide_idx)


    # for information of the render status
    def get_render_status(self) -> dict:
        rendered = sum(1 for s in (self.slide_cache or []) if s != b"")
        return {
            "total":    self.total_slides,
            "rendered": rendered,
            "ready":    rendered == self.total_slides,
        }


    # converts the PPTX/PPT into PDF first, since libreoffice directly converting to images will not work or causing troubles
    # using subprocess, make sure libreoffice is installed
    # if not, the rendered PPT will only be text-only, no fancy graphics like images
    def _convert_to_pdf(self, filepath: str) -> tuple:
        try:
            tmpdir = tempfile.mkdtemp(prefix="pptprocessor_")
            result = subprocess.run(
                [
                    "libreoffice", "--headless",
                    f"-env:UserInstallation=file://{tmpdir}/lo_profile",
                    "--convert-to", "pdf",
                    "--outdir", tmpdir,
                    filepath,
                ],
                capture_output=True,
                timeout=120,
            )
 
            print(f"[LibreOffice] returncode: {result.returncode}")
            if result.stderr:
                print(f"[LibreOffice] stderr: {result.stderr.decode()}")
 
            if result.returncode != 0:
                shutil.rmtree(tmpdir, ignore_errors=True)
                return None, None
 
            from pathlib import Path
            pdf_path = os.path.join(tmpdir, Path(filepath).stem + ".pdf")
            if not os.path.exists(pdf_path):
                print(f"[LibreOffice] PDF not found at: {pdf_path}")
                shutil.rmtree(tmpdir, ignore_errors=True)
                return None, None
 
            print(f"[LibreOffice] PDF ready: {pdf_path}")
            return pdf_path, tmpdir
 
        except Exception as exc:
            print(f"[PPTProcessor] PDF conversion error: {exc}")
            return None, None

    # using pdf2image to convert each of the presentation slides into JPEG image format
    # using the pdf2image library installed from the pip
    def _render_via_pdf(self, slide_idx: int) -> Optional[bytes]:
        # Render a slide by extracting the corresponding page from the PDF.
        from pdf2image import convert_from_path
 
        if not self.pdf_path or not os.path.exists(self.pdf_path):
            return None
 
        try:
            images = convert_from_path(
                self.pdf_path,
                dpi=150,
                first_page=slide_idx + 1,
                last_page=slide_idx + 1,
            )
            if not images:
                return None
 
            buf = io.BytesIO()
            images[0].save(buf, format="JPEG", quality=RENDER_QUALITY)
            return buf.getvalue()
 
        except Exception as exc:
            print(f"[PPTProcessor] pdf2image render error (slide {slide_idx}): {exc}")
            return None

    # this using python-pptx will only render text and no images
    # this will likely to be called if the libreoffice not failed to converts the file
    def _render_python_fallback(self, prs: Presentation, slide_idx: int) -> bytes:
        print(f"[PPTProcessor] Rendering slide {slide_idx} via Python fallback")
        slide = prs.slides[slide_idx]

        w_emu: int = 0
        h_emu: int = 0

        try:
            w_emu: int = int(prs.slide_width  or 9144000)
            h_emu: int = int(prs.slide_height or 5143500)

            # 914400 is the EMUs per inch2
            scale = DEFAULT_WIDTH / (w_emu / 914400 * 96)
            width = DEFAULT_WIDTH
            height = int((h_emu / 914400 * 96) * scale)
        except Exception:
            width, height = DEFAULT_WIDTH, DEFAULT_HEIGHT

        bg_color = (255,255,255)
        try:
            bg = slide.background
            fill = bg.fill
            if fill.type is not None:
                from pptx.util import Pt
                from pptx.dml.color import RGBColor

                fg = fill.fore_color
                rgb = fg.rgb
                bg_color = (rgb[0], rgb[1], rgb[2])
        except Exception:
            pass

        img = Image.new("RGB", (width, height), bg_color)

        try:
            draw = ImageDraw.Draw(img)
            try:
                font_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36)
                font_body = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)

            except Exception:
                font_title = ImageFont.load_default()
                font_body = ImageFont.load_default()

            for base_shape in slide.shapes:
                if not base_shape.has_text_frame:
                    continue

                try:
                    x = int(base_shape.left / w_emu * width)
                    y = int(base_shape.top / h_emu * height)
                    sw = int(base_shape.width / w_emu * width)
                except Exception:
                    x, y, sw = 20, 20, width - 40

                cursor_y = y

                shape = cast(Shape, base_shape)
                text_frame = shape.text_frame
                for para_idx, para in enumerate(text_frame.paragraphs):
                    line = para.text.strip()

                    if not line:
                        cursor_y += 12
                        continue

                    font = font_title if para_idx == 0 else font_body
                    color = (30,30,30) # dark grey

                    try:
                        run = para.runs[0]
                        rgb = run.font.color.rgb
                        color = tuple(rgb) # red. green, blue
                    except Exception:
                        pass

                    draw.text((x + 8, cursor_y), line, fill=color, font=font)
                    try:
                        bbox = font.getbbox(line)
                        line_height = (bbox[3] - bbox[1]) + 6
                    except Exception:
                        line_height = 30

                    cursor_y += line_height
        except Exception as exc:
            print(f"[PPTProcessor] Text render error on slide {slide_idx}: {exc}")

        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=RENDER_QUALITY)


        print(f"[PPTProcessor] Slide {slide_idx} rendered, size: {len(buf.getvalue())} bytes")
        return buf.getvalue()
